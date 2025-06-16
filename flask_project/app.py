import os
import re
import json
import tempfile
from flask import Flask, request, jsonify, send_file
from pptx import Presentation
from pptx.util import Inches
from openai import OpenAI

# === SETUP ===
app = Flask(__name__)
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

ppt_data = {
    "presentation": None,
    "file_path": None
}

# === CORE UTILITIES ===
def load_pptx(file):
    return Presentation(file)

def save_pptx(prs):
    path = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx").name
    prs.save(path)
    return path

def get_slide_text(slide):
    return "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))

def insert_placeholder_icon(slide, left=Inches(5), top=Inches(1)):
    placeholder_path = os.path.join(os.path.dirname(__file__), "placeholder_icon.png")
    if os.path.exists(placeholder_path):
        slide.shapes.add_picture(placeholder_path, left, top, height=Inches(1))

# === ACTION HANDLERS ===
def handle_edit_text(prs, args):
    slides_to_modify = prs.slides
    if isinstance(args["slide_number"], int):
        slides_to_modify = [prs.slides[args["slide_number"] - 1]]

    for slide in slides_to_modify:
        for shape in slide.shapes:
            if hasattr(shape, "text") and args["old_text"] in shape.text:
                shape.text = shape.text.replace(args["old_text"], args["new_text"])
                return "Edit applied."
    return "Text not found."


def handle_apply_style(prs, args):
    slide = prs.slides[args["slide_number"] - 1]
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            shape.text_frame.paragraphs[0].font.name = "Arial"
            shape.text_frame.paragraphs[0].font.size = Inches(0.4)
    return "Modern style applied."

def handle_add_icon(prs, args):
    slide = prs.slides[args["slide_number"] - 1]
    insert_placeholder_icon(slide)
    return "Placeholder icon added."

def handle_change_layout(prs, args):
    slide = prs.slides[args["slide_number"] - 1]
    shapes = slide.shapes
    if len(shapes) >= 1 and hasattr(shapes[0], "text"):
        shapes[0].left = Inches(0.5)
        shapes[0].top = Inches(0.5)
        if len(shapes) > 1:
            shapes[1].left = Inches(4)
            shapes[1].top = Inches(0.5)
    return "Slide layout updated."

# === GPT PARSER ===
def parse_chat_and_apply(prs, user_msg):
    slides_info = [
        {"slide_number": i + 1, "text": get_slide_text(slide)}
        for i, slide in enumerate(prs.slides)
    ]

    messages = [
        {
            "role": "system",
            "content": (
                "You are a PowerPoint assistant. Based on the user's instruction and slide content, return only JSON.\n"
                "Examples:\n"
                "{\"action\": \"edit_text\", \"slide_number\": 1, \"old_text\": \"Welcome\", \"new_text\": \"Quarterly Review\"}\n"
                "{\"action\": \"apply_style\", \"slide_number\": 2, \"style\": \"modern\"}\n"
                "{\"action\": \"add_icon\", \"slide_number\": 3, \"keywords\": [\"profit\"]}\n"
            )
        },
        {
            "role": "user",
            "content": f"Slides: {slides_info}\nInstruction: {user_msg}"
        }
    ]

    response = client.chat.completions.create(
        model="gpt-4",
        messages=messages
    )

    raw = response.choices[0].message.content
    print("GPT RAW RESPONSE:\n", raw)

    try:
        json_text = re.search(r'\{[\s\S]*\}', raw).group()
        result = json.loads(json_text)
    except Exception as e:
        return f"Failed to parse GPT response: {str(e)}"

    action_map = {
        "edit_text": handle_edit_text,
        "apply_style": handle_apply_style,
        "add_icon": handle_add_icon,
        "change_layout": handle_change_layout
    }

    action_func = action_map.get(result["action"])
    return action_func(prs, result) if action_func else "Unknown action."

# === FLASK ROUTES ===
@app.route("/upload", methods=["POST"])
def upload():
    file = request.files['file']
    prs = load_pptx(file)
    ppt_data["presentation"] = prs
    ppt_data["file_path"] = save_pptx(prs)
    return jsonify({"message": "Presentation loaded."})

@app.route("/chat", methods=["POST"])
def chat():
    if ppt_data["presentation"] is None:
        return jsonify({"error": "No presentation loaded."}), 400

    user_msg = request.json.get("message", "")
    print(f"Received user message: {user_msg}")

    try:
        response = parse_chat_and_apply(ppt_data["presentation"], user_msg)
        ppt_data["file_path"] = save_pptx(ppt_data["presentation"])
        return jsonify({"response": response})
    except Exception as e:
        print("Error:", e)
        return jsonify({"error": str(e)}), 500

@app.route("/download", methods=["GET"])
def download():
    if ppt_data["file_path"]:
        return send_file(ppt_data["file_path"], as_attachment=True)
    return jsonify({"error": "No file available."})

if __name__ == "__main__":
    app.run(debug=True, port=5000)